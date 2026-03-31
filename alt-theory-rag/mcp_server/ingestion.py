"""Document Ingestion System for Knowledge RAG

Multi-format document parsing, chunking, and metadata extraction.
Supports: MD, PDF, TXT, PY, JSON, DOCX, XLSX, PPTX, CSV
"""

import hashlib
import json
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

# PDF support (optional)
try:
    import fitz  # PyMuPDF

    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

# Office formats (optional)
try:
    import docx  # python-docx

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl

    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

import csv
import io

from .config import config


@dataclass
class Chunk:
    """A chunk of text from a document"""

    content: str
    index: int
    start_char: int
    end_char: int
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Document:
    """Parsed document with metadata and chunks"""

    id: str
    content: str
    source: Path
    format: str
    category: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    chunks: List[Chunk] = field(default_factory=list)
    keywords: List[str] = field(default_factory=list)

    @property
    def filename(self) -> str:
        return self.source.name

    @property
    def relative_path(self) -> str:
        try:
            return str(self.source.relative_to(config.documents_dir))
        except ValueError:
            return str(self.source)


class DocumentParser:
    """Multi-format document parser with chunking and metadata extraction"""

    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or config.chunk_size
        self.chunk_overlap = chunk_overlap or config.chunk_overlap

        # Parser dispatch table
        self._parsers = {
            ".md": self._parse_markdown,
            ".txt": self._parse_text,
            ".pdf": self._parse_pdf,
            ".py": self._parse_code,
            ".json": self._parse_json,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".csv": self._parse_csv,
        }

    def parse_file(self, filepath: Path) -> Optional[Document]:
        """Parse a file and return a Document object with chunks"""
        filepath = Path(filepath)

        if not filepath.exists():
            raise FileNotFoundError(f"File not found: {filepath}")

        suffix = filepath.suffix.lower()
        if suffix not in self._parsers:
            raise ValueError(f"Unsupported format: {suffix}")

        # Generate unique ID
        doc_id = self._generate_id(filepath)

        # Parse content and metadata
        content, metadata = self._parsers[suffix](filepath)

        if not content or not content.strip():
            return None

        # Detect category from path
        category = self._detect_category(filepath)

        # Extract keywords
        keywords = self._extract_keywords(content, category)

        # Create document
        doc = Document(
            id=doc_id,
            content=content,
            source=filepath,
            format=suffix,
            category=category,
            metadata=metadata,
            keywords=keywords,
        )

        # Chunk the content (markdown-aware for .md files)
        if suffix == ".md":
            doc.chunks = self._chunk_markdown(content, metadata)
        else:
            doc.chunks = self._chunk_text(content, metadata)

        return doc

    def parse_directory(self, directory: Path = None) -> List[Document]:
        """Parse all supported files in a directory recursively"""
        directory = Path(directory) if directory else config.documents_dir
        documents = []

        for ext in config.supported_formats:
            for filepath in directory.rglob(f"*{ext}"):
                try:
                    doc = self.parse_file(filepath)
                    if doc:
                        documents.append(doc)
                except Exception as e:
                    print(f"[WARN] Failed to parse {filepath}: {e}")

        return documents

    # =========================================================================
    # Format-specific parsers
    # =========================================================================

    def _parse_markdown(self, filepath: Path) -> tuple[str, Dict]:
        """Parse Markdown file, extracting headers as metadata"""
        content = filepath.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "type": "markdown",
            "headers": [],
            "has_code_blocks": "```" in content,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
        }

        # Extract headers hierarchy
        header_pattern = r"^(#{1,6})\s+(.+)$"
        for match in re.finditer(header_pattern, content, re.MULTILINE):
            level = len(match.group(1))
            title = match.group(2).strip()
            metadata["headers"].append({"level": level, "title": title})

        # Extract title from first H1 or filename
        h1_headers = [h for h in metadata["headers"] if h["level"] == 1]
        if h1_headers:
            metadata["title"] = h1_headers[0]["title"]
        else:
            metadata["title"] = filepath.stem

        # Extract frontmatter if present (YAML between ---)
        frontmatter_match = re.match(r"^---\n(.*?)\n---\n", content, re.DOTALL)
        if frontmatter_match:
            metadata["has_frontmatter"] = True
            # Remove frontmatter from content for cleaner indexing
            content = content[frontmatter_match.end() :]

        return content, metadata

    def _parse_pdf(self, filepath: Path) -> tuple[str, Dict]:
        """Parse PDF file using PyMuPDF (text extraction, no markdown conversion)."""
        if not HAS_PYMUPDF:
            raise ImportError("PyMuPDF (fitz) not installed. Install with: pip install pymupdf")

        metadata = {
            "type": "pdf",
            "pages": 0,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
        }

        text_parts = []

        with fitz.open(filepath) as doc:
            metadata["pages"] = len(doc)
            metadata["title"] = doc.metadata.get("title", filepath.stem)
            metadata["author"] = doc.metadata.get("author", "")

            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")

        content = "\n\n".join(text_parts)
        return content, metadata

    def _parse_text(self, filepath: Path) -> tuple[str, Dict]:
        """Parse plain text file"""
        content = filepath.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "type": "text",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
            "line_count": content.count("\n") + 1,
        }
        return content, metadata

    def _parse_code(self, filepath: Path) -> tuple[str, Dict]:
        """Parse Python code file, extracting docstrings and comments"""
        content = filepath.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "type": "code",
            "language": "python",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
            "functions": [],
            "classes": [],
            "imports": [],
        }

        # Extract module docstring
        docstring_match = re.match(r'^["\'][\'"]{2}(.*?)["\'][\'"]{2}', content, re.DOTALL)
        if docstring_match:
            metadata["docstring"] = docstring_match.group(1).strip()

        # Extract function names
        func_pattern = r"^def\s+(\w+)\s*\("
        metadata["functions"] = re.findall(func_pattern, content, re.MULTILINE)

        # Extract class names
        class_pattern = r"^class\s+(\w+)\s*[:\(]"
        metadata["classes"] = re.findall(class_pattern, content, re.MULTILINE)

        # Extract imports
        import_pattern = r"^(?:from\s+[\w.]+\s+)?import\s+[\w.,\s]+"
        metadata["imports"] = re.findall(import_pattern, content, re.MULTILINE)[:10]

        return content, metadata

    def _parse_json(self, filepath: Path) -> tuple[str, Dict]:
        """Parse JSON file"""
        raw_content = filepath.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "type": "json",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
        }

        try:
            data = json.loads(raw_content)
            metadata["is_valid_json"] = True

            if isinstance(data, dict):
                metadata["keys"] = list(data.keys())[:20]
                metadata["structure"] = "object"
            elif isinstance(data, list):
                metadata["length"] = len(data)
                metadata["structure"] = "array"

            # Pretty-print for better indexing
            content = json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            metadata["is_valid_json"] = False
            content = raw_content

        return content, metadata

    def _parse_docx(self, filepath: Path) -> tuple[str, Dict]:
        """Parse DOCX file extracting paragraphs and tables."""
        if not HAS_DOCX:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

        doc = docx.Document(filepath)
        metadata = {
            "type": "docx",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
        }

        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure as markdown
                if para.style and para.style.name.startswith("Heading"):
                    try:
                        level = int(para.style.name.split()[-1])
                        parts.append(f"{'#' * level} {text}")
                    except (ValueError, IndexError):
                        parts.append(f"## {text}")
                else:
                    parts.append(text)

        # Extract tables as markdown
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

        content = "\n\n".join(parts)
        return content, metadata

    def _parse_xlsx(self, filepath: Path) -> tuple[str, Dict]:
        """Parse XLSX file extracting all sheets as text tables."""
        if not HAS_XLSX:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")

        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        metadata = {
            "type": "xlsx",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
            "sheets": wb.sheetnames,
        }

        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != " | " * (len(cells) - 1):
                    parts.append(line)

        wb.close()
        content = "\n\n".join(parts)
        return content, metadata

    def _parse_pptx(self, filepath: Path) -> tuple[str, Dict]:
        """Parse PPTX file extracting slide text."""
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        prs = Presentation(filepath)
        metadata = {
            "type": "pptx",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
            "slides": len(prs.slides),
        }

        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"## Slide {i + 1}\n" + "\n".join(slide_texts))

        content = "\n\n".join(parts)
        return content, metadata

    def _parse_csv(self, filepath: Path) -> tuple[str, Dict]:
        """Parse CSV file as text table."""
        raw = filepath.read_text(encoding="utf-8", errors="ignore")
        metadata = {
            "type": "csv",
            "title": filepath.stem,
            "file_size": filepath.stat().st_size,
            "modified": datetime.fromtimestamp(filepath.stat().st_mtime).isoformat(),
        }

        parts = []
        reader = csv.reader(io.StringIO(raw))
        rows = list(reader)
        metadata["rows"] = len(rows)
        metadata["columns"] = len(rows[0]) if rows else 0

        for row in rows:
            parts.append(" | ".join(row))

        content = "\n".join(parts)
        return content, metadata

    # =========================================================================
    # Chunking
    # =========================================================================

    def _chunk_text(self, text: str, metadata: Dict) -> List[Chunk]:
        """Split text into overlapping chunks for embedding"""
        if not text:
            return []

        chunks = []
        text_len = len(text)
        start = 0
        index = 0
        previous_start = -1  # Track previous start to detect infinite loops

        while start < text_len:
            # Safety: detect infinite loop (start not progressing)
            if start <= previous_start:
                break
            previous_start = start

            # Calculate end position
            end = min(start + self.chunk_size, text_len)

            # Try to break at sentence/paragraph boundary
            if end < text_len:
                # Look for natural break points within last 20% of chunk
                break_zone_start = start + int(self.chunk_size * 0.8)
                break_zone = text[break_zone_start:end]

                # Priority: paragraph > sentence > word
                for pattern in ["\n\n", "\n", ". ", " "]:
                    last_break = break_zone.rfind(pattern)
                    if last_break != -1:
                        end = break_zone_start + last_break + len(pattern)
                        break

            chunk_content = text[start:end].strip()

            if chunk_content:
                chunk = Chunk(
                    content=chunk_content,
                    index=index,
                    start_char=start,
                    end_char=end,
                    metadata={
                        "title": metadata.get("title", ""),
                        "type": metadata.get("type", ""),
                    },
                )
                chunks.append(chunk)
                index += 1

            # Move start position with overlap
            # Ensure we always make forward progress
            new_start = end - self.chunk_overlap

            # If overlap would cause no progress, just move to end
            if new_start <= start:
                start = end
            else:
                start = new_start

        return chunks

    def _chunk_markdown(self, text: str, metadata: Dict) -> List[Chunk]:
        """
        Markdown-aware chunking with code block protection and min-size merging.

        1. Strips code blocks before splitting (prevents # comments from being treated as headers)
        2. Splits by ## and ### headers only (not # which catches code comments)
        3. Merges small chunks (<min_chunk_size) with the next section
        4. Falls back to _chunk_text() if no headers found

        Args:
            text: Full document text
            metadata: Document metadata dict

        Returns:
            List of Chunk objects aligned to markdown sections
        """
        if not text:
            return []

        min_chunk_size = 100  # Minimum chars for a standalone chunk

        # Step 1: Mask code blocks to prevent splitting on # inside them
        code_blocks = []

        def mask_code(match):
            code_blocks.append(match.group(0))
            return f"__CODE_BLOCK_{len(code_blocks) - 1}__"

        masked_text = re.sub(r"```.*?```", mask_code, text, flags=re.DOTALL)

        # Step 2: Split by ## and ### headers only (not # which catches code comments)
        sections = re.split(r"(?=^#{2,3}\s+)", masked_text, flags=re.MULTILINE)

        # Filter empty sections
        sections = [s for s in sections if s.strip()]

        if len(sections) <= 1:
            return self._chunk_text(text, metadata)

        # Step 3: Restore code blocks in each section
        def restore_code(section_text):
            for i, block in enumerate(code_blocks):
                section_text = section_text.replace(f"__CODE_BLOCK_{i}__", block)
            return section_text

        sections = [restore_code(s) for s in sections]

        # Step 4: Merge small sections with the next one
        merged_sections = []
        buffer = ""
        for section in sections:
            if buffer:
                buffer += "\n\n" + section
                if len(buffer.strip()) >= min_chunk_size:
                    merged_sections.append(buffer)
                    buffer = ""
            elif len(section.strip()) < min_chunk_size:
                buffer = section
            else:
                merged_sections.append(section)

        if buffer:
            if merged_sections:
                merged_sections[-1] += "\n\n" + buffer
            else:
                merged_sections.append(buffer)

        if not merged_sections:
            return self._chunk_text(text, metadata)

        # Step 5: Create chunks from merged sections
        chunks = []
        global_index = 0
        char_offset = 0

        for section in merged_sections:
            section_stripped = section.strip()
            if not section_stripped:
                char_offset += len(section)
                continue

            header_match = re.match(r"^(#{2,3}\s+.+)$", section_stripped, re.MULTILINE)
            header_context = header_match.group(1) if header_match else ""

            if len(section_stripped) <= self.chunk_size:
                chunk = Chunk(
                    content=section_stripped,
                    index=global_index,
                    start_char=char_offset,
                    end_char=char_offset + len(section),
                    metadata={
                        "title": metadata.get("title", ""),
                        "type": metadata.get("type", ""),
                        "section_header": header_context,
                    },
                )
                chunks.append(chunk)
                global_index += 1
            else:
                sub_chunks = self._chunk_text(section_stripped, metadata)
                for i, sub_chunk in enumerate(sub_chunks):
                    if i > 0 and header_context:
                        sub_chunk.content = f"{header_context}\n\n{sub_chunk.content}"
                    sub_chunk.index = global_index
                    sub_chunk.start_char += char_offset
                    sub_chunk.end_char += char_offset
                    sub_chunk.metadata["section_header"] = header_context
                    chunks.append(sub_chunk)
                    global_index += 1

            char_offset += len(section)

        return chunks

    # =========================================================================
    # Category detection
    # =========================================================================

    def _detect_category(self, filepath: Path) -> str:
        """Detect document category based on file path"""
        try:
            rel_path = filepath.relative_to(config.documents_dir)
            path_str = str(rel_path).replace("\\", "/").lower()
        except ValueError:
            path_str = str(filepath).replace("\\", "/").lower()

        # Check category mappings in order (more specific first)
        for path_pattern, category in sorted(config.category_mappings.items(), key=lambda x: len(x[0]), reverse=True):
            if path_pattern in path_str:
                return category

        return "general"

    # =========================================================================
    # Keyword extraction
    # =========================================================================

    def _extract_keywords(self, content: str, category: str) -> List[str]:
        """Extract technical keywords from content"""
        keywords = set()
        content_lower = content.lower()

        # Check against all keyword routes
        for route_category, route_keywords in config.keyword_routes.items():
            for keyword in route_keywords:
                if keyword.lower() in content_lower:
                    keywords.add(keyword.lower())

        # Extract additional technical terms
        # CVE patterns
        cve_pattern = r"CVE-\d{4}-\d{4,}"
        keywords.update(re.findall(cve_pattern, content, re.IGNORECASE))

        # MITRE ATT&CK patterns
        mitre_pattern = r"T\d{4}(?:\.\d{3})?"
        keywords.update(re.findall(mitre_pattern, content))

        # IP addresses
        ip_pattern = r"\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b"
        ips = re.findall(ip_pattern, content)
        if len(ips) <= 5:  # Only add if not too many (likely real targets)
            keywords.update(ips)

        # Common security tools mentioned
        security_tools = [
            "nmap",
            "burp",
            "metasploit",
            "wireshark",
            "hydra",
            "john",
            "hashcat",
            "gobuster",
            "nikto",
            "sqlmap",
            "nuclei",
            "ffuf",
            "bloodhound",
            "mimikatz",
            "responder",
            "crackmapexec",
            "impacket",
        ]
        for tool in security_tools:
            if tool in content_lower:
                keywords.add(tool)

        return sorted(list(keywords))

    # =========================================================================
    # Utilities
    # =========================================================================

    def _generate_id(self, filepath: Path) -> str:
        """Generate unique document ID based on path and modification time"""
        stat = filepath.stat()
        unique_str = f"{filepath}:{stat.st_mtime}:{stat.st_size}"
        return hashlib.sha256(unique_str.encode()).hexdigest()[:16]


# Convenience function
def parse_documents(directory: Path = None) -> List[Document]:
    """Parse all documents in a directory"""
    parser = DocumentParser()
    return parser.parse_directory(directory)
